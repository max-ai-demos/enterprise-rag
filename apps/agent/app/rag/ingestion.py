# apps/agent/app/rag/ingestion.py
#
# PDF 解析默认使用 PyMuPDF（fitz），适用于文字版 PDF。
#
# 商业化升级路径（扫描版 OCR、精确表格/图片解析）：
#   将 .env 中的 PDF_PARSER_PROVIDER 改为 "datalab" 或 "mineru"，并填写对应 Key/URL。
#   - datalab_client.py / mineru_client.py 已就绪，切换后会自动走对应解析路径。
#   - Datalab 返回带精确 bbox 的结构化 JSON，可同时填充 document_index 表（存图片 URL）。
#   - 如需持久化解析结果，文件存储从本地磁盘升级为 S3，缓存解析结果避免重复计算。
#
import logging
from pathlib import Path
from typing import Any
from sqlalchemy.orm import Session
from llama_index.core import Settings as LlamaSettings
from llama_index.core.node_parser import SentenceSplitter
from llama_index.embeddings.openai import OpenAIEmbedding
from app.infrastructure.config import settings

logger = logging.getLogger(__name__)

CHUNK_SIZE = 600
CHUNK_OVERLAP = 80


def _get_splitter():
    """Return the configured chunker.

    CHUNK_STRATEGY=sentence (default): LlamaIndex SentenceSplitter
    CHUNK_STRATEGY=smart:              段落优先、句子兜底的 SmartChunker
    切换策略后需重新运行 reingest_all.py。
    """
    if settings.chunk_strategy == "smart":
        from app.rag.chunking import SmartChunker
        return SmartChunker(
            target_size=settings.smart_chunk_target_size,
            min_size=settings.smart_chunk_min_size,
        )
    return SentenceSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)

_WORD_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def _parse_docx_blocks(file_path: str) -> list[dict]:
    """Extract paragraphs and table rows from DOCX in document order."""
    from docx import Document as DocxDocument
    from docx.table import Table as DocxTable

    doc = DocxDocument(file_path)
    blocks = []
    idx = 0

    for child in doc.element.body:
        tag = child.tag
        if tag == f"{_WORD_NS}p":
            text = "".join(child.itertext()).strip()
            if text:
                blocks.append({"text": text, "block_idx": idx, "is_table": False})
                idx += 1
        elif tag == f"{_WORD_NS}tbl":
            tbl = DocxTable(child, doc)
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    row_text = " | ".join(cells)
                    blocks.append({"text": row_text, "block_idx": idx, "is_table": True})
                    idx += 1
    return blocks


def _setup_llama_settings():
    LlamaSettings.embed_model = OpenAIEmbedding(
        model="text-embedding-3-small",
        api_key=settings.openai_api_key,
    )


def parse_document(file_path: str, file_type: str) -> list[dict[str, Any]]:
    """Parse a document into chunks with position metadata."""
    chunks = []

    if file_type == "pdf":
        import fitz
        doc = fitz.open(file_path)
        chunk_index = 0
        splitter = _get_splitter()
        for page_num in range(len(doc)):
            page = doc[page_num]
            pw = page.rect.width or 1
            ph = page.rect.height or 1
            blocks = page.get_text("blocks")
            for block in blocks:
                x0, y0, x1, y1, text, _block_no, block_type = block
                if block_type != 0 or not text.strip():
                    continue
                # Normalize bbox to [0..1000] for frontend compatibility
                bbox = [
                    round(x0 / pw * 1000, 1),
                    round(y0 / ph * 1000, 1),
                    round(x1 / pw * 1000, 1),
                    round(y1 / ph * 1000, 1),
                ]
                sub_texts = splitter.split_text(text.strip()) if len(text.strip()) > CHUNK_SIZE else [text.strip()]
                for sub in sub_texts:
                    if sub.strip():
                        chunks.append({
                            "text": sub.strip(),
                            "page_num": page_num + 1,
                            "page_idx": page_num + 1,
                            "bbox": bbox,
                            "chunk_index": chunk_index,
                        })
                        chunk_index += 1

    elif file_type in ("ppt", "pptx"):
        from pptx import Presentation
        prs = Presentation(file_path)
        splitter = _get_splitter()
        chunk_index = 0
        for slide_idx, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if not texts:
                continue
            slide_text = "\n".join(texts)
            sub_texts = splitter.split_text(slide_text) if hasattr(splitter, "split_text") else [slide_text]
            for sub in sub_texts:
                if sub.strip():
                    chunks.append({
                        "text": sub.strip(),
                        "slide_num": slide_idx + 1,
                        "chunk_index": chunk_index,
                    })
                    chunk_index += 1

    elif file_type == "docx":
        splitter = _get_splitter()
        chunk_index = 0
        for block in _parse_docx_blocks(file_path):
            for sub in splitter.split_text(block["text"]):
                if sub.strip():
                    chunks.append({
                        "text": sub.strip(),
                        "paragraph_idx": block["block_idx"],
                        "chunk_index": chunk_index,
                    })
                    chunk_index += 1

    elif file_type == "xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        chunk_index = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                text = " | ".join(str(c) for c in row if c is not None)
                if text.strip():
                    chunks.append({
                        "text": text.strip(),
                        "sheet_name": sheet_name,
                        "row_start": row_idx + 1,
                        "row_end": row_idx + 1,
                        "chunk_index": chunk_index,
                    })
                    chunk_index += 1

    elif file_type == "txt":
        text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        splitter = _get_splitter()
        for chunk_index, sub in enumerate(splitter.split_text(text)):
            if sub.strip():
                chunks.append({
                    "text": sub.strip(),
                    "char_offset": 0,
                    "chunk_index": chunk_index,
                })

    return chunks


def _parse_pdf_cloud(file_path: str) -> list[dict]:
    """Parse PDF via Datalab or MinerU and return chunks in the same format as parse_document."""
    import asyncio
    import json as _json

    provider = settings.pdf_parser_provider
    file_content = Path(file_path).read_bytes()
    file_name = Path(file_path).name
    splitter = SentenceSplitter(chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    chunks = []

    if provider == "datalab":
        from app.infrastructure.datalab_client import DatalabClient
        client = DatalabClient()
        _markdown, extracted_dir = asyncio.run(client.process_pdf(file_content, file_name))
        pdf_name = Path(file_path).stem
        content_list_path = extracted_dir / pdf_name / f"{pdf_name}_content_list.json"
        content_list = _json.loads(content_list_path.read_text(encoding="utf-8"))
        chunk_index = 0
        for item in content_list:
            text = item.get("text") or item.get("table_body") or ""
            if not text.strip():
                continue
            page_idx = item.get("page_idx", 0)
            bbox = item.get("bbox") or []
            for sub in splitter.split_text(text.strip()):
                if sub.strip():
                    chunks.append({
                        "text": sub.strip(),
                        "page_num": page_idx + 1,
                        "page_idx": page_idx + 1,
                        "bbox": bbox,
                        "chunk_index": chunk_index,
                    })
                    chunk_index += 1

    elif provider == "mineru":
        from app.infrastructure.mineru_client import MineruClient
        client = MineruClient()
        markdown = client.process_pdf(file_content, file_name)
        chunk_index = 0
        for sub in splitter.split_text(markdown):
            if sub.strip():
                chunks.append({
                    "text": sub.strip(),
                    "page_num": 0,
                    "page_idx": 0,
                    "bbox": [],
                    "chunk_index": chunk_index,
                })
                chunk_index += 1

    else:
        logger.warning(f"Unknown pdf_parser_provider '{provider}', falling back to PyMuPDF")
        return parse_document(file_path, "pdf")

    logger.info(f"Cloud PDF parse ({provider}): {len(chunks)} chunks from {file_name}")
    return chunks


def ingest_document(document_id: str, file_path: str, file_type: str, db: Session) -> int:
    """Parse, embed, and store chunks in MySQL. Returns chunk count.

    PDF parsing provider is controlled by PDF_PARSER_PROVIDER in .env:
      "pymupdf"  — default, no extra config
      "datalab"  — requires DATALAB_API_KEY
      "mineru"   — requires MINERU_BASE_URL
    """
    from app.db.repository import ChunkRepository

    _setup_llama_settings()

    # Use cloud parser when configured; fall back to PyMuPDF otherwise.
    if file_type == "pdf" and settings.pdf_parser_provider != "pymupdf":
        chunks = _parse_pdf_cloud(file_path)
    else:
        chunks = parse_document(file_path, file_type)
    if not chunks:
        logger.warning(f"No chunks from {file_path}")
        return 0

    chunk_repo = ChunkRepository(db)
    embed_model = LlamaSettings.embed_model
    batch_size = 100

    for i in range(0, len(chunks), batch_size):
        batch = chunks[i:i + batch_size]
        embeddings = embed_model.get_text_embedding_batch([c["text"] for c in batch])
        chunk_repo.upsert_batch(document_id, batch, embeddings, file_type)

    return len(chunks)


def delete_document_vectors(document_id: str, db: Session):
    from app.db.repository import ChunkRepository
    ChunkRepository(db).delete_by_document(document_id)
